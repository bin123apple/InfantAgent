"""
agentskills.py

This module provides various file manipulation skills for the Infant agent.

"""
import os
import re
import time
import base64
import shutil
import chardet
import tempfile
import functools
import mimetypes
import subprocess
import pandas as pd 
from math import sqrt
from inspect import signature
from typing import Optional

import docx
import pypdf 
from openai import OpenAI
from pptx import Presentation
from pylatexenc.latex2text import LatexNodes2Text
from PIL import ImageGrab, ImageDraw, ImageFont, Image

CURRENT_FILE: str | None = None
CURRENT_LINE = 1
WINDOW = 100

ENABLE_AUTO_LINT = os.getenv('ENABLE_AUTO_LINT', 'false').lower() == 'true'
EXCEL_EXTENSIONS = {".xls", ".xlsx"}
# This is also used in unit tests!
MSG_FILE_UPDATED = '[File updated. Please review the changes and make sure they are correct (correct indentation, no duplicate lines, etc). Edit the file again if necessary.]'

# OPENAI
OPENAI_API_KEY = os.getenv(
    'OPENAI_API_KEY', os.getenv('SANDBOX_ENV_OPENAI_API_KEY', '')
)
OPENAI_BASE_URL = os.getenv('OPENAI_BASE_URL', 'https://api.openai.com/v1')
OPENAI_MODEL = os.getenv('OPENAI_MODEL', 'gpt-4o-2024-05-13')
MAX_TOKEN = os.getenv('MAX_TOKEN', 500)

OPENAI_PROXY = f'{OPENAI_BASE_URL}/chat/completions'

client = OpenAI(api_key=OPENAI_API_KEY, base_url=OPENAI_BASE_URL)

### File Operations ###

# Define the decorator using the functionality of UpdatePwd
# FIXME: Do we need to check the cd command? As it will change the path.
def update_pwd_decorator(func):
    @functools.wraps(func)
    def wrapper(*args, **kwargs):
        old_pwd = os.getcwd()
        jupyter_pwd = os.environ.get('JUPYTER_PWD', None)
        if jupyter_pwd:
            os.chdir(jupyter_pwd)
        try:
            return func(*args, **kwargs)
        finally:
            os.chdir(old_pwd)

    return wrapper

def _is_valid_filename(file_name) -> bool:
    if not file_name or not isinstance(file_name, str) or not file_name.strip():
        return False
    invalid_chars = '<>:"/\\|?*'
    if os.name == 'nt':  # Windows
        invalid_chars = '<>:"/\\|?*'
    elif os.name == 'posix':  # Unix-like systems
        invalid_chars = '\0'

    for char in invalid_chars:
        if char in file_name:
            return False
    return True

def _is_valid_path(path) -> bool:
    if not path or not isinstance(path, str):
        return False
    try:
        return os.path.exists(os.path.normpath(path))
    except PermissionError:
        return False

def _create_paths(file_name) -> bool:
    try:
        dirname = os.path.dirname(file_name)
        if dirname:
            os.makedirs(dirname, exist_ok=True)
        return True
    except PermissionError:
        return False

def _check_current_file(file_path: str | None = None) -> bool:
    global CURRENT_FILE
    if not file_path:
        file_path = CURRENT_FILE
    if not file_path or not os.path.isfile(file_path):
        raise ValueError('No file open. Use the open_file function first.')
    return True

def _clamp(value, min_value, max_value):
    return max(min_value, min(value, max_value))

def match_str(my_string, my_list, line_number):
    '''
    Match a string from a list.
    If no match, try to use the fuzzy matching mode
    '''
    closest_index = None
    closest_distance = float('inf')
    
    if my_string == '':
        return closest_index
    
    for i, element in enumerate(my_list):
        if my_string == element.replace('\n',''):
            distance = abs(i + 1 - line_number)
            if distance < closest_distance:
                closest_index = i
                closest_distance = distance

    if closest_index is None:
        for i, element in enumerate(my_list):
            if my_string.strip() in element.strip():
                distance = abs(i + 1 - line_number)
                if distance < closest_distance:
                    closest_index = i
                    closest_distance = distance

    if closest_index is not None:
        return closest_index
    else:
        return None

def _lint_file(file_path: str) -> tuple[Optional[str], Optional[int]]:
    """
    Lint the file at the given path and return a tuple with a boolean indicating if there are errors,
    and the line number of the first error, if any.

    Returns:
        tuple[str, Optional[int]]: (lint_error, first_error_line_number)
    """

    if file_path.endswith('.py'):
        # Define the flake8 command with selected error codes
        command = [
            'flake8',
            '--isolated',
            '--select=F821,F822,F831,E112,E113,E999,E902',
            file_path,
        ]

        # Run the command using subprocess and redirect stderr to stdout
        result = subprocess.run(
            command,
            stdout=subprocess.PIPE,
            stderr=subprocess.STDOUT,
        )
        if result.returncode == 0:
            # Linting successful. No issues found.
            return None, None

        # Extract the line number from the first error message
        error_message = result.stdout.decode().strip()
        lint_error = 'ERRORS:\n' + error_message
        first_error_line = None
        for line in error_message.splitlines(True):
            if line.strip():
                # The format of the error message is: <filename>:<line>:<column>: <error code> <error message>
                parts = line.split(':')
                if len(parts) >= 2:
                    try:
                        first_error_line = int(parts[1])
                        break
                    except ValueError:
                        # Not a valid line number, continue to the next line
                        continue

        return lint_error, first_error_line

    # Not a python file, skip linting
    return None, None

def _cur_file_header(CURRENT_FILE, total_lines) -> str:
    if not CURRENT_FILE:
        return ''
    return f'[File: {os.path.abspath(CURRENT_FILE)} ({total_lines} lines total)]\n'

def search_function_line_number(file_path: str, function_signature: str) -> tuple[int | None, int | None]:
    """
    Args:
        file_path (str): The path of the file to search.
        function_signature (str): The function name to search for.
        
    Returns:
        tuple[int | None, int | None]: A tuple (start_line, end_line) indicating the lines where the function is defined.
            If not found, returns (None, None) and prints an error message.
    """
    try:
        with open(file_path, 'r') as file:
            code_lines = file.readlines()
        
        start_line_number = None
        end_line_number = None
        header_end_line = None 
        indent_level = None
        
        pattern = re.compile(rf'^\s*def\s+{re.escape(function_signature)}\b')
        
        for index, line in enumerate(code_lines):
            if pattern.match(line):
                start_line_number = index + 1  
                header_lines = [line]
                header_end_line = index
                if not line.rstrip().endswith(':'):
                    for j in range(index + 1, len(code_lines)):
                        header_lines.append(code_lines[j])
                        if code_lines[j].rstrip().endswith(':'):
                            header_end_line = j
                            break
                indent_level = len(code_lines[index]) - len(code_lines[index].lstrip())
                
                for k in range(header_end_line + 1, len(code_lines)):
                    current_line = code_lines[k]
                    if current_line.strip() == '':
                        continue
                    current_indent = len(current_line) - len(current_line.lstrip())
                    if current_indent <= indent_level:
                        end_line_number = k
                        break
                if end_line_number is None:
                    end_line_number = len(code_lines)
                return (start_line_number, end_line_number)
        
        print(f"Function '{function_signature}' not found in the file '{file_path}'.")
        print("Please check the function signature and try again. NOTE: The function signature should only include the function name.")
        return None, None
    except FileNotFoundError:
        raise Exception(f"File '{file_path}' not found.")
    except Exception as e:
        raise Exception(f"An error occurred: {e}")

def _clamp(value, min_value, max_value):
    return max(min(value, max_value), min_value)

def _cur_file_header(file_path, total_lines):
    return f"File: {file_path} (Total lines: {total_lines})\n"

def _print_window(file_path, targeted_line, WINDOW, return_str=False):
    global CURRENT_LINE
    _check_current_file(file_path)
    with open(file_path) as file:
        content = file.read()

        # Ensure the content ends with a newline character
        if not content.endswith('\n'):
            content += '\n'

        lines = content.splitlines(True)  # Keep all line ending characters
        total_lines = len(lines)

        # cover edge cases
        CURRENT_LINE = _clamp(targeted_line, 1, total_lines)
        half_window = max(1, WINDOW // 2)

        # Ensure at least one line above and below the targeted line
        start = max(1, CURRENT_LINE - half_window)
        end = min(total_lines, CURRENT_LINE + half_window)

        # Adjust start and end to ensure at least one line above and below
        if start == 1:
            end = min(total_lines, start + WINDOW - 1)
        if end == total_lines:
            start = max(1, end - WINDOW + 1)

        output = ''

        # only display this when there's at least one line above
        if start > 1:
            output += f'({start - 1} more lines above)\n'
        for i in range(start, end + 1):
            _new_line = f'{i}|{lines[i-1]}'
            if not _new_line.endswith('\n'):
                _new_line += '\n'
            output += _new_line
        if end < total_lines:
            output += f'({total_lines - end} more lines below)\n'
        output = output.rstrip()

        if return_str:
            return output
        else:
            print(output)

def is_text_file(file_path):
    """
    Determine if a file is a text file based on its content and extension.
    """
    try:
        mime_type, _ = mimetypes.guess_type(file_path)
        if mime_type and mime_type.startswith("text"):
            return True

        with open(file_path, "rb") as f:
            raw_data = f.read(2048)
            if b"\x00" in raw_data: 
                return False
            result = chardet.detect(raw_data)
            encoding = result["encoding"]

        return bool(encoding)
    except Exception as e:
        print(f"Error detecting file type: {e}")
        return False  

@update_pwd_decorator
def open_file(
    path: str, line_number: int | None = 1, context_lines: int | None = 150
) -> None:
    """
    Opens the file at the given path in the editor. If line_number is provided, the window will be moved to include that line.
    It only shows the first 100 lines by default! Max `context_lines` supported is 2000, use `scroll up/down`
    to view the file if you want to see more.

    Args:
        path: str: The path to the file to open, preferredly absolute path.
        line_number: int | None = 1: The line number to move to. Defaults to 1.
        context_lines: int | None = 100: Only shows this number of lines in the context window (usually from line 1), with line_number as the center (if possible). Defaults to 100.
    """
    global CURRENT_FILE, CURRENT_LINE, WINDOW

    if not os.path.isfile(path):
        raise FileNotFoundError(f"File {path} not found")

    _, ext = os.path.splitext(path)
    ext = ext.lower()
    if ext in EXCEL_EXTENSIONS:
        CURRENT_FILE = os.path.abspath(path)
        df = pd.read_excel(CURRENT_FILE)
        print(f"Opening Excel file: {CURRENT_FILE}")
        print(f"Displaying first {context_lines} rows:")
        print(df.head(context_lines))
        return
    elif not is_text_file(path):
        raise ValueError(f"Unsupported binary file: {ext}. Only text files are supported.")

    CURRENT_FILE = os.path.abspath(path)
    with open(CURRENT_FILE, encoding="utf-8") as file:
        total_lines = max(1, sum(1 for _ in file))

    if not isinstance(line_number, int) or line_number < 1 or line_number > total_lines:
        raise ValueError(f"Line number must be between 1 and {total_lines}")
    CURRENT_LINE = line_number

    # Override WINDOW with context_lines
    if context_lines is None or context_lines < 1:
        context_lines = 150
    WINDOW = _clamp(context_lines, 1, 2000)

    output = _cur_file_header(CURRENT_FILE, total_lines)
    output += _print_window(CURRENT_FILE, CURRENT_LINE, WINDOW, return_str=True)
    print(output)

@update_pwd_decorator
def goto_line(line_number: int) -> None:
    """
    Moves the window to show the specified line number.

    Args:
        line_number: int: The line number to move to.
    """
    global CURRENT_FILE, CURRENT_LINE, WINDOW
    _check_current_file()

    with open(str(CURRENT_FILE)) as file:
        total_lines = max(1, sum(1 for _ in file))
    if not isinstance(line_number, int) or line_number < 1 or line_number > total_lines:
        raise ValueError(f'Line number must be between 1 and {total_lines}')

    CURRENT_LINE = _clamp(line_number, 1, total_lines)

    output = _cur_file_header(CURRENT_FILE, total_lines)
    output += _print_window(CURRENT_FILE, CURRENT_LINE, WINDOW, return_str=True)
    print(output)

@update_pwd_decorator
def scroll_down() -> None:
    """Moves the window down by 100 lines.

    Args:
        None
    """
    global CURRENT_FILE, CURRENT_LINE, WINDOW
    _check_current_file()

    with open(str(CURRENT_FILE)) as file:
        total_lines = max(1, sum(1 for _ in file))
    CURRENT_LINE = _clamp(CURRENT_LINE + WINDOW, 1, total_lines)
    output = _cur_file_header(CURRENT_FILE, total_lines)
    output += _print_window(CURRENT_FILE, CURRENT_LINE, WINDOW, return_str=True)
    print(output)

@update_pwd_decorator
def scroll_up() -> None:
    """Moves the window up by 100 lines.

    Args:
        None
    """
    global CURRENT_FILE, CURRENT_LINE, WINDOW
    _check_current_file()

    with open(str(CURRENT_FILE)) as file:
        total_lines = max(1, sum(1 for _ in file))
    CURRENT_LINE = _clamp(CURRENT_LINE - WINDOW, 1, total_lines)
    output = _cur_file_header(CURRENT_FILE, total_lines)
    output += _print_window(CURRENT_FILE, CURRENT_LINE, WINDOW, return_str=True)
    print(output)

@update_pwd_decorator
def create_file(filename: str) -> None:
    """Creates and opens a new file with the given name.

    Args:
        filename: str: The name of the file to create.
    """
    if os.path.exists(filename):
        raise FileExistsError(f"File '{filename}' already exists.")

    with open(filename, 'w') as file:
        file.write('\n')

    open_file(filename)
    print(f'[File {filename} created.]')
    
@update_pwd_decorator
def search_function(file_path: str, function_signature: str) -> str:
    """
    Args:
        file_path (str): The path of the file to search.
        function_signature (str): The signature of the function to search for.

    Returns:
        str: The code of the function if found, or an error message if not found.
    """
    try:
        if not file_path.endswith('.py'):
            raise ValueError(f"search_function only supports python file, but got {file_path}")
        start_line_number, end_line_number = search_function_line_number(file_path, function_signature)
        if start_line_number is None or end_line_number is None:
            print(f"Can not match '{function_signature}'. Possible reasons:\n"
                  f"1. The code style of file '{file_path}' is not supported, or\n"
                  f"2. The function signature is incorrect (please provide only the function name).")
            return f"search_function does not support this file type or function definition style."
        context_lines = end_line_number - start_line_number + 1
        line_number = (end_line_number + start_line_number) // 2
        open_file(path = file_path, line_number = line_number, context_lines = context_lines)
    except Exception as e:
        raise e

@update_pwd_decorator
def replace_function(
    file_name: str, 
    function_to_replace: str,
    new_code: str
    ) -> None:
    """Replace some lines inside the function 
    (This is used to avoid some potential wrong line numberproblems in edit_file funtion)

    Args:
        filename: str: The name of the file to create.
        code_to_replace: str: The origianl code that will be replaced
        new_code: str: new code that will be used
    """
    global CURRENT_FILE, CURRENT_LINE, WINDOW

    # Check file
    
    if not _is_valid_filename(file_name):
        
        raise FileNotFoundError('Invalid file name.')

    if not _is_valid_path(file_name):
        
        raise FileNotFoundError('Invalid path or file name.')

    if not _create_paths(file_name):
        
        raise PermissionError('Could not access or create directories.')

    if not os.path.isfile(file_name):
        
        raise FileNotFoundError(f'File {file_name} not found.')

    # Use a temporary file to write changes
    temp_file_path = ''
    src_abs_path = os.path.abspath(file_name)
    first_error_line = None
    try:
        # # lint the origianl file
        if ENABLE_AUTO_LINT:
            origianl_lint_error, origianl_first_error_line = _lint_file(file_name)      
              
        with tempfile.NamedTemporaryFile('w', delete=False) as temp_file:
            temp_file_path = temp_file.name

            # Read the original file and check if empty and for a trailing newline
            with open(file_name) as original_file:
                file_content = original_file.read()   
                start_line, end_line = search_function_line_number(file_name, function_to_replace)
                if start_line is None or end_line is None:
                    return
                def replace_code_block(file_content, start_line, end_line, new_code):
                    file_lines = file_content.splitlines()
                    file_lines[start_line-1:end_line] = [new_code]
                    updated_file_content = "\n".join(file_lines)
                    return updated_file_content
                
                # replace the code block
                updated_file_content = replace_code_block(file_content, start_line, end_line, new_code)
                
            # Write the new content to the temporary file
            temp_file.write(updated_file_content)

        # Replace the original file with the temporary file atomically
        shutil.move(temp_file_path, src_abs_path)
        
        # find the first different line
        original_lines = file_content.splitlines()
        updated_lines = updated_file_content.splitlines()
        
        new_code_line = None
        for index, (original, updated) in enumerate(zip(original_lines, updated_lines), start=1):
            if original != updated:
                new_code_line = index
                break
                
        # Modified lines
        m_lines = len(new_code.splitlines())
        
        # middle screen
        window = m_lines + 10
        # print(f"window: {window}")
        if new_code_line is not None:
            middle_screen = new_code_line
        else:
            middle_screen = None
            print("The new code is exactly the same as the original code. You should use a different code and try again!")
        # print(f"middle_screen: {middle_screen}")
        
        
        # Handle linting
        if ENABLE_AUTO_LINT:
            # BACKUP the original file
            original_file_backup_path = os.path.join(
                os.path.dirname(file_name),
                f'.backup.{os.path.basename(file_name)}',
            )
            with open(original_file_backup_path, 'w') as f:
                 f.write(file_content)

            lint_error, first_error_line = _lint_file(file_name)

            # Select the errors caused by the modification
            def extract_last_part(line):
                parts = line.split(':')
                if len(parts) > 1:
                    return parts[-1].strip()
                return line.strip()

            def subtract_strings(str1, str2) -> str:
                lines1 = str1.splitlines()
                lines2 = str2.splitlines()
                
                last_parts1 = [extract_last_part(line) for line in lines1]
                
                remaining_lines = [line for line in lines2 if extract_last_part(line) not in last_parts1]

                result = '\n'.join(remaining_lines)
                return result
            if origianl_lint_error:
                lint_error = subtract_strings(origianl_lint_error, lint_error)
                if lint_error == "":
                    lint_error = None
                    first_error_line = None
                                                           
            if lint_error is not None:
                
                if first_error_line is not None:
                    CURRENT_LINE = middle_screen
                print(
                    '[Your proposed edit has introduced new syntax error(s). Please understand the errors and retry your edit command.]'
                )
                
                print('[This is how your edit would have looked if applied]')
                print('-------------------------------------------------')
                _print_window(file_name, CURRENT_LINE, window)
                print('-------------------------------------------------\n')

                print('[This is the original code before your edit]')
                print('-------------------------------------------------')
                _print_window(original_file_backup_path, CURRENT_LINE, window)
                print('-------------------------------------------------')

                print(
                    'Your changes have NOT been applied. Please fix your edit command and try again based on the following error messages.\n'
                    f'{lint_error}\n'
                    'You may need to do one or several of the following:\n'
                    '1) Specify the correct code block that you want to modify;\n' 
                    '2) Make sure that the Args position is correct (the 2nd arg is the origianl code that will be replaced and the 3rd arg is the new code that will be used);\n' 
                    '3) Choose another command (Such as edit_file(file_name, start, start_str end, end_str, content) command).\n'
                    '4) Use open_file(path, line_number, context_lines) command to check the details of where you want to modify and improve your command.\n'
                    'DO NOT re-run the same failed edit command. Running it again will lead to the same error.'
                )

                # recover the original file
                with open(original_file_backup_path) as fin, open(
                    file_name, 'w'
                ) as fout:
                    fout.write(fin.read())
                os.remove(original_file_backup_path)
                return

    except FileNotFoundError as e:
        
        print(f'File not found: {e}')
    except IOError as e:
        
        print(f'An error occurred while handling the file: {e}')
    except ValueError as e:
        
        print(f'Invalid input: {e}')
    except Exception as e:
        
        # Clean up the temporary file if an error occurs
        if temp_file_path and os.path.exists(temp_file_path):
            os.remove(temp_file_path)
        print(f'An unexpected error occurred: {e}')
        raise e

    # Update the file information and print the updated content
    with open(file_name, 'r', encoding='utf-8') as file:
        n_total_lines = max(1, len(file.readlines()))
    if first_error_line is not None and int(first_error_line) > 0:
        CURRENT_LINE = first_error_line
    else:
        CURRENT_LINE = middle_screen or n_total_lines or 1
    print(
        f'[File: {os.path.abspath(file_name)} ({n_total_lines} lines total after edit)]'
    )
    CURRENT_FILE = file_name
    _print_window(CURRENT_FILE, CURRENT_LINE, WINDOW)
    print(MSG_FILE_UPDATED)

@update_pwd_decorator
def replace_code(
    file_name: str, 
    code_to_replace: str,
    new_code: str
    ) -> None:
    """Replace some lines inside the function 
    (This is used to avoid some potential wrong line numberproblems in edit_file funtion)

    Args:
        filename: str: The name of the file to create.
        code_to_replace: str: The origianl code that will be replaced
        new_code: str: new code that will be used
    """
    global CURRENT_FILE, CURRENT_LINE, WINDOW

    # Check file
    
    if not _is_valid_filename(file_name):
        
        raise FileNotFoundError('Invalid file name.')

    if not _is_valid_path(file_name):
        
        raise FileNotFoundError('Invalid path or file name.')

    if not _create_paths(file_name):
        
        raise PermissionError('Could not access or create directories.')

    if not os.path.isfile(file_name):
        
        raise FileNotFoundError(f'File {file_name} not found.')

    # Use a temporary file to write changes
    temp_file_path = ''
    src_abs_path = os.path.abspath(file_name)
    first_error_line = None
    try:
        # # lint the origianl file
        if ENABLE_AUTO_LINT:
            origianl_lint_error, origianl_first_error_line = _lint_file(file_name)      
              
        with tempfile.NamedTemporaryFile('w', delete=False) as temp_file:
            temp_file_path = temp_file.name

            # Read the original file and check if empty and for a trailing newline
            with open(file_name) as original_file:
                file_content = original_file.read()   
                def flexible_str(s: str) -> str:
                    s = re.escape(s)
                    flexible_pattern = re.sub(r'\n', r'\n+', s) # change \n to \n\n...\n
                    return flexible_pattern
                
                flexible_code_to_replace = flexible_str(code_to_replace)
                matchs = re.findall(flexible_code_to_replace, file_content)
                occurrences = len(matchs)
                
                # arguments check
                if occurrences == 0:
                    print(
                        f'The code block:\n'
                        f'{code_to_replace}\n'
                        f'is not involved in the {file_name}.\n'
                        'Your changes have NOT been applied.\n'
                        'Please use open_file(path, line_number, context_lines) command to check the details of where you want to modify and fix your command.\n'
                        'Or You can also use the edit_file(file_name, start, start_str end, end_str, content) command to indicate the code block that you want to modify.'
                    )
                    return

                if occurrences > 1:
                    print(
                        f'The code block:\n'
                        f'{code_to_replace}\n'
                        f'is duplicated in the {file_name}.\n'
                        'Your changes have NOT been applied.\n'
                        'Please use the edit_file(file_name, start, start_str end, end_str, content) command to indicate the code block that you want to modify.'
                    )
                    return
                
                # replace the code block
                code_to_replace = matchs[0]
                updated_file_content = file_content.replace(code_to_replace, new_code)
                
            # Write the new content to the temporary file
            temp_file.write(updated_file_content)

        # Replace the original file with the temporary file atomically
        shutil.move(temp_file_path, src_abs_path)
        
        # find the first different line
        original_lines = file_content.splitlines()
        updated_lines = updated_file_content.splitlines()

        for index, (original, updated) in enumerate(zip(original_lines, updated_lines), start=1):
            if original != updated:
                new_code_line = index
                break
                
        # Modified lines
        m_lines = len(new_code.splitlines())
        
        # middle screen
        window = m_lines + 10
        # print(f"window: {window}")
        middle_screen = new_code_line
        # print(f"middle_screen: {middle_screen}")
        
        
        # Handle linting
        if ENABLE_AUTO_LINT:
            # BACKUP the original file
            original_file_backup_path = os.path.join(
                os.path.dirname(file_name),
                f'.backup.{os.path.basename(file_name)}',
            )
            with open(original_file_backup_path, 'w') as f:
                 f.write(file_content)

            lint_error, first_error_line = _lint_file(file_name)

            # Select the errors caused by the modification
            def extract_last_part(line):
                parts = line.split(':')
                if len(parts) > 1:
                    return parts[-1].strip()
                return line.strip()

            def subtract_strings(str1, str2) -> str:
                lines1 = str1.splitlines()
                lines2 = str2.splitlines()
                
                last_parts1 = [extract_last_part(line) for line in lines1]
                
                remaining_lines = [line for line in lines2 if extract_last_part(line) not in last_parts1]

                result = '\n'.join(remaining_lines)
                return result
            if origianl_lint_error:
                lint_error = subtract_strings(origianl_lint_error, lint_error)
                if lint_error == "":
                    lint_error = None
                    first_error_line = None
                                                           
            if lint_error is not None:
                
                if first_error_line is not None:
                    CURRENT_LINE = middle_screen
                print(
                    '[Your proposed edit has introduced new syntax error(s). Please understand the errors and retry your edit command.]'
                )
                
                print('[This is how your edit would have looked if applied]')
                print('-------------------------------------------------')
                _print_window(file_name, CURRENT_LINE, window)
                print('-------------------------------------------------\n')

                print('[This is the original code before your edit]')
                print('-------------------------------------------------')
                _print_window(original_file_backup_path, CURRENT_LINE, window)
                print('-------------------------------------------------')

                print(
                    'Your changes have NOT been applied. Please fix your edit command and try again based on the following error messages.\n'
                    f'{lint_error}\n'
                    'You may need to do one or several of the following:\n'
                    '1) Specify the correct code block that you want to modify;\n' 
                    '2) Make sure that the Args position is correct (the 2nd arg is the origianl code that will be replaced and the 3rd arg is the new code that will be used);\n' 
                    '3) Choose another command (Such as edit_file(file_name, start, start_str end, end_str, content) command).\n'
                    '4) Use open_file(path, line_number, context_lines) command to check the details of where you want to modify and improve your command.\n'
                    'DO NOT re-run the same failed edit command. Running it again will lead to the same error.'
                )

                # recover the original file
                with open(original_file_backup_path) as fin, open(
                    file_name, 'w'
                ) as fout:
                    fout.write(fin.read())
                os.remove(original_file_backup_path)
                return

    except FileNotFoundError as e:
        
        print(f'File not found: {e}')
    except IOError as e:
        
        print(f'An error occurred while handling the file: {e}')
    except ValueError as e:
        
        print(f'Invalid input: {e}')
    except Exception as e:
        
        # Clean up the temporary file if an error occurs
        if temp_file_path and os.path.exists(temp_file_path):
            os.remove(temp_file_path)
        print(f'An unexpected error occurred: {e}')
        raise e

    # Update the file information and print the updated content
    with open(file_name, 'r', encoding='utf-8') as file:
        n_total_lines = max(1, len(file.readlines()))
    if first_error_line is not None and int(first_error_line) > 0:
        CURRENT_LINE = first_error_line
    else:
        CURRENT_LINE = middle_screen or n_total_lines or 1
    print(
        f'[File: {os.path.abspath(file_name)} ({n_total_lines} lines total after edit)]'
    )
    CURRENT_FILE = file_name
    _print_window(CURRENT_FILE, CURRENT_LINE, WINDOW)
    print(MSG_FILE_UPDATED)

def _edit_or_append_file(
    file_name: str,
    start: int | None = None,
    start_str: str | None = None,
    end: int | None = None,
    end_str: str| None = None,
    content: str = '',
    is_append: bool = False,
) -> None:
    """Internal method to handle common logic for edit_/append_file methods.

    Args:
        file_name: str: The name of the file to edit or append to.
        start: int | None = None: The start line number for editing. Ignored if is_append is True.
        end: int | None = None: The end line number for editing. Ignored if is_append is True.
        content: str: The content to replace the lines with or to append.
        is_append: bool = False: Whether to append content to the file instead of editing.
    """
    global CURRENT_FILE, CURRENT_LINE, WINDOW

    ERROR_MSG = f'[Error editing file {file_name}. Please confirm the file is correct.]'
    ERROR_MSG_SUFFIX = (
        'Your changes have NOT been applied. Please fix your edit command and try again.\n'
        'You either need to 1) Open the correct file and try again or 2) Specify the correct start/end line arguments.\n'
        'DO NOT re-run the same failed edit command. Running it again will lead to the same error.'
    )
    if not _is_valid_filename(file_name):
        
        raise FileNotFoundError('Invalid file name.')

    if not _is_valid_path(file_name):
        
        raise FileNotFoundError('Invalid path or file name.')

    if not _create_paths(file_name):
        
        raise PermissionError('Could not access or create directories.')

    if not os.path.isfile(file_name):
        
        raise FileNotFoundError(f'File {file_name} not found.')

    # Use a temporary file to write changes
    content = str(content or '')
    temp_file_path = ''
    src_abs_path = os.path.abspath(file_name)
    first_error_line = None
    try:
        # # lint the origianl file
        if ENABLE_AUTO_LINT:
            origianl_lint_error, origianl_first_error_line = _lint_file(file_name)
        # Create a temporary file
        with tempfile.NamedTemporaryFile('w', delete=False) as temp_file:
            temp_file_path = temp_file.name
            # Read the original file and check if empty and for a trailing newline
            with open(file_name) as original_file:
                lines = original_file.readlines()

            if is_append:
                if lines and not (len(lines) == 1 and lines[0].strip() == ''):
                    if not lines[-1].endswith('\n'):
                        lines[-1] += '\n'
                    content_lines = content.splitlines(keepends=True)
                    new_lines = lines + content_lines
                    content = ''.join(new_lines)
            else:
                # Handle cases where start or end are None
                if start is None:
                    start = 1  # Default to the beginning
                if end is None:
                    end = len(lines)  # Default to the end
                            
                # Check arguments
                if not (1 <= start <= len(lines)):
                    
                    print(
                        f'{ERROR_MSG}\n'
                        f'Invalid start line number: {start}. Line numbers must be between 1 and {len(lines)} (inclusive).\n'
                        f'{ERROR_MSG_SUFFIX}'
                    )
                    return
                if not (1 <= end <= len(lines)):
                    
                    print(
                        f'{ERROR_MSG}\n'
                        f'Invalid end line number: {end}. Line numbers must be between 1 and {len(lines)} (inclusive).\n'
                        f'{ERROR_MSG_SUFFIX}'
                    )
                    return
                if start > end:
                    
                    print(
                        f'{ERROR_MSG}\n'
                        f'Invalid line range: {start}-{end}. Start must be less than or equal to end.\n'
                        f'{ERROR_MSG_SUFFIX}'
                    )
                    return
                
                # Double check if the start/end line is correct by compare with the start_str/end_str
                check_pass = True
                if lines[start-1].replace('\n','') != start_str:
                    print(f'The string: {start_str} does not match the start line: {start}')
                    # print(f'The start line: {start} is {lines[start-1]}')
                    check_pass = False
                if lines[end-1].replace('\n','') != end_str:
                    print(f'The string: {end_str} does not match the end line: {end}')    
                    # print(f'The end line: {end} is {lines[end-1]}')   
                    check_pass = False    
                if check_pass != True:
                    print("Here is the code that you are trying to modified:")
                    window_size = (end - start) + 5
                    target_line = (end + start) // 2
                    start_str_closest_index = match_str(start_str, lines, start)
                    end_str_closest_index = match_str(end_str, lines, end)
                    _print_window(file_name, target_line, window_size)
                    print(
                        f'The start line: {start} is:\n'
                        f'{start}|{lines[start-1]}\n'
                        f'The end line: {end} is:\n'
                        f'{end}|{lines[end-1]}\n'
                    )
                    if start_str_closest_index is not None:
                        print(
                            f"The matching string closest to the line {start} and most similar to the start_str you provided is at position {start_str_closest_index + 1}.\n"
                            f'{start_str_closest_index + 1}|{lines[start_str_closest_index]}'
                        )
                    if end_str_closest_index is not None:
                        print(
                            f"The matching string closest to the line {end} and most similar to the end_str you provided is at position {end_str_closest_index + 1}.\n"
                            f'{end_str_closest_index + 1}|{lines[end_str_closest_index]}'
                        )                    
                    print(
                        'Your changes have NOT been applied. Please fix your edit command and try again.\n'
                        'Please double-check whether this part of the code is what you originally planned to modify\n'
                        "If you want to use the edit_file() command, please provide the correct start line and end line along with the corresponding strings on those lines. And don't forget to provide the `content` argument.\n"
                        "You should first try to use the information above to modify your edit_file() command.\n"
                        'However, if you have already tried to fix this edit_file() command multiple times and the same issue persists, please try using replace_function() to modify the code.'
                    )
                    return

                if not content.endswith('\n'):
                    content += '\n'
                content_lines = content.splitlines(True)
                new_lines = lines[: start - 1] + content_lines + lines[end:]
                content = ''.join(new_lines)

            if not content.endswith('\n'):
                content += '\n'

            # Write the new content to the temporary file
            temp_file.write(content)

        # Replace the original file with the temporary file atomically
        shutil.move(temp_file_path, src_abs_path)

        # Handle linting
        if ENABLE_AUTO_LINT:
            # BACKUP the original file
            original_file_backup_path = os.path.join(
                os.path.dirname(file_name),
                f'.backup.{os.path.basename(file_name)}',
            )
            with open(original_file_backup_path, 'w') as f:
                f.writelines(lines)

            lint_error, first_error_line = _lint_file(file_name)
            
            # Select the errors caused by the modification
            
            def extract_last_part(line):
                parts = line.split(':')
                if len(parts) > 1:
                    return parts[-1].strip()
                return line.strip()

            def subtract_strings(str1, str2) -> str:
                lines1 = str1.splitlines()
                lines2 = str2.splitlines()
                
                last_parts1 = [extract_last_part(line) for line in lines1]
                
                remaining_lines = [line for line in lines2 if extract_last_part(line) not in last_parts1]

                result = '\n'.join(remaining_lines)
                return result
            if origianl_lint_error:
                lint_error = subtract_strings(origianl_lint_error, lint_error)
                if lint_error == "":
                    lint_error = None
                    first_error_line = None
                            
            if lint_error is not None:
                
                if first_error_line is not None:
                    # CURRENT_LINE = int(first_error_line)
                    total_lines = len(content_lines)
                    CURRENT_LINE = int(start) + total_lines // 2 # Try the edited line
                print(
                    '[Your proposed edit has introduced new syntax error(s). Please understand the errors and retry your edit command.]'
                )
                
                print('[This is how your edit would have looked if applied]')
                print('-------------------------------------------------')
                _print_window(file_name, CURRENT_LINE, total_lines + 10)
                print('-------------------------------------------------\n')

                print('[This is the original code before your edit]')
                print('-------------------------------------------------')
                _print_window(original_file_backup_path, CURRENT_LINE, total_lines + 10)
                print('-------------------------------------------------')

                print(
                    'Your changes have NOT been applied. Please fix your edit command and try again based on the following error messages.\n'
                    f'{lint_error}\n'
                    'You probably need to do one or several of the following:\n'
                    '1) Specify the correct start/end line parameters;\n' 
                    '2) Correct your edit code;\n' 
                    '3) Choose another command (Such as replace_function(file_name,code_to_replace,new_code) command).\n'
                    '4) Use open_file(path, line_number, context_lines) command to check the details of where you want to modify and improve your command\n'
                    'DO NOT re-run the same failed edit command. Running it again will lead to the same error.'
                )

                # recover the original file
                with open(original_file_backup_path) as fin, open(
                    file_name, 'w'
                ) as fout:
                    fout.write(fin.read())
                os.remove(original_file_backup_path)
                return

    except FileNotFoundError as e:
        
        print(f'File not found: {e}')
    except IOError as e:
        
        print(f'An error occurred while handling the file: {e}')
    except ValueError as e:
        
        print(f'Invalid input: {e}')
    except Exception as e:
        
        # Clean up the temporary file if an error occurs
        if temp_file_path and os.path.exists(temp_file_path):
            os.remove(temp_file_path)
        print(f'An unexpected error occurred: {e}')
        raise e

    # Update the file information and print the updated content
    with open(file_name, 'r', encoding='utf-8') as file:
        n_total_lines = max(1, len(file.readlines()))
    if first_error_line is not None and int(first_error_line) > 0:
        CURRENT_LINE = first_error_line
    else:
        if is_append:
            CURRENT_LINE = max(1, len(lines))  # end of original file
        else:
            CURRENT_LINE = start or n_total_lines or 1
    print(
        f'[File: {os.path.abspath(file_name)} ({n_total_lines} lines total after edit)]'
    )
    CURRENT_FILE = file_name
    _print_window(CURRENT_FILE, CURRENT_LINE, WINDOW)
    print(MSG_FILE_UPDATED)

@update_pwd_decorator
def edit_file(file_name: str, start: int, start_str: str, end: int, end_str: str, content: str) -> None:
    """Edit a file.

    Replaces in given file `file_name` the lines `start` through `end` (inclusive) with the given text `content`.
    If a line must be inserted, an already existing line must be passed in `content` with new content accordingly!
    The parameters start_str and end_str will be used to verify if the lines you specified are correct. 
    start_str should correspond to the string on the start line, and end_str should correspond to the string on the end line.
    To avoid potential errors, the string corresponding to content should first be assigned to a variable, and then that variable should be used in the command.
    
    For example:
    EDIT_CODE = '''
    <CONTENT>
    '''
    edit_file(file_name, start, start_str, end, end_str, content=EDIT_CODE)
    
    Args:
        file_name: str: The name of the file to edit.
        start: int: The start line number. Must satisfy start >= 1.
        start_str: str: String on the start line.
        end: int: The end line number. Must satisfy start <= end <= number of lines in the file.
        end_str: str: String on the end line.
        content: str: The content to replace the lines with.
    """
    _edit_or_append_file(
        file_name, start=start, start_str=start_str, end=end, end_str=end_str, content=content, is_append=False
    )

@update_pwd_decorator
def append_file(file_name: str, content: str) -> None:
    """Append content to the given file.

    It appends text `content` to the end of the specified file.

    Args:
        file_name: str: The name of the file to append to.
        content: str: The content to append to the file.
    """
    _edit_or_append_file(file_name, start=1, end=None, content=content, is_append=True)

@update_pwd_decorator
def search_dir(search_term: str, dir_path: str = './') -> None:
    """Searches for search_term in all files in dir. If dir is not provided, searches in the current directory.

    Args:
        search_term: str: The term to search for.
        dir_path: Optional[str]: The path to the directory to search.
    """
    if not os.path.isdir(dir_path):
        raise FileNotFoundError(f'Directory {dir_path} not found')
    matches = []
    for root, _, files in os.walk(dir_path):
        for file in files:
            if file.startswith('.'):
                continue
            file_path = os.path.join(root, file)
            with open(file_path, 'r', errors='ignore') as f:
                for line_num, line in enumerate(f, 1):
                    if search_term in line:
                        matches.append((file_path, line_num, line.strip()))

    if not matches:
        print(f'No matches found for "{search_term}" in {dir_path}')
        return

    num_matches = len(matches)
    num_files = len(set(match[0] for match in matches))

    if num_files > 100:
        print(
            f'More than {num_files} files matched for "{search_term}" in {dir_path}. Please narrow your search.'
        )
        return

    print(f'[Found {num_matches} matches for "{search_term}" in {dir_path}]')
    for file_path, line_num, line in matches:
        print(f'{file_path} (Line {line_num}): {line}')
    print(f'[End of matches for "{search_term}" in {dir_path}]')

@update_pwd_decorator
def search_file(search_term: str, file_path: Optional[str] = None) -> None:
    """Searches for search_term in file. If file is not provided, searches in the current open file.

    Args:
        search_term: str: The term to search for.
        file_path: Optional[str]: The path to the file to search.
    """
    global CURRENT_FILE
    if file_path is None:
        file_path = CURRENT_FILE
    if file_path is None:
        raise FileNotFoundError(
            'No file specified or open. Use the open_file function first.'
        )
    if not os.path.isfile(file_path):
        raise FileNotFoundError(f'File {file_path} not found')

    matches = []
    with open(file_path) as file:
        for i, line in enumerate(file, 1):
            if search_term in line:
                matches.append((i, line.strip()))

    if matches:
        print(f'[Found {len(matches)} matches for "{search_term}" in {file_path}]')
        for match in matches:
            print(f'Line {match[0]}: {match[1]}')
        print(f'[End of matches for "{search_term}" in {file_path}]')
    else:
        print(f'[No matches found for "{search_term}" in {file_path}]')

@update_pwd_decorator
def find_file(file_name: str, dir_path: str = './') -> None:
    """Finds all files with the given name in the specified directory.

    Args:
        file_name: str: The name of the file to find.
        dir_path: Optional[str]: The path to the directory to search.
    """
    if not os.path.isdir(dir_path):
        raise FileNotFoundError(f'Directory {dir_path} not found')

    matches = []
    for root, _, files in os.walk(dir_path):
        for file in files:
            if file_name in file:
                matches.append(os.path.join(root, file))

    if matches:
        print(f'[Found {len(matches)} matches for "{file_name}" in {dir_path}]')
        for match in matches:
            print(f'{match}')
        print(f'[End of matches for "{file_name}" in {dir_path}]')
    else:
        print(f'[No matches found for "{file_name}" in {dir_path}]')

### Analysis different kinds of files ###

@update_pwd_decorator
def parse_pdf(file_path: str) -> None:
    """Parses the content of a PDF file and prints it.

    Args:
        file_path: str: The path to the file to open.
    """
    print(f'[Reading PDF file from {file_path}]')
    content = pypdf.PdfReader(file_path)
    text = ''
    for page_idx in range(len(content.pages)):
        text += (
            f'@@ Page {page_idx + 1} @@\n'
            + content.pages[page_idx].extract_text()
            + '\n\n'
        )
    print(text.strip())

@update_pwd_decorator
def parse_docx(file_path: str) -> None:
    """
    Parses the content of a DOCX file and prints it.

    Args:
        file_path: str: The path to the file to open.
    """
    print(f'[Reading DOCX file from {file_path}]')
    content = docx.Document(file_path)
    text = ''
    for i, para in enumerate(content.paragraphs):
        text += f'@@ Page {i + 1} @@\n' + para.text + '\n\n'
    print(text)

@update_pwd_decorator
def parse_latex(file_path: str) -> None:
    """
    Parses the content of a LaTex file and prints it.

    Args:
        file_path: str: The path to the file to open.
    """
    print(f'[Reading LaTex file from {file_path}]')
    with open(file_path) as f:
        data = f.read()
    text = LatexNodes2Text().latex_to_text(data)
    print(text.strip())

def _base64_img(file_path: str) -> str:
    with open(file_path, 'rb') as image_file:
        encoded_image = base64.b64encode(image_file.read()).decode('utf-8')
    return encoded_image

def _base64_video(file_path: str, frame_interval: int = 10) -> list[str]:
    import cv2

    video = cv2.VideoCapture(file_path)
    base64_frames = []
    frame_count = 0
    while video.isOpened():
        success, frame = video.read()
        if not success:
            break
        if frame_count % frame_interval == 0:
            _, buffer = cv2.imencode('.jpg', frame)
            base64_frames.append(base64.b64encode(buffer).decode('utf-8'))
        frame_count += 1
    video.release()
    return base64_frames

def _prepare_image_messages(task: str, base64_image: str):
    return [
        {
            'role': 'user',
            'content': [
                {'type': 'text', 'text': task},
                {
                    'type': 'image_url',
                    'image_url': {'url': f'data:image/jpeg;base64,{base64_image}'},
                },
            ],
        }
    ]

@update_pwd_decorator
def parse_audio(file_path: str, model: str = 'whisper-1') -> None:
    """
    Parses the content of an audio file and prints it.

    Args:
        file_path: str: The path to the audio file to transcribe.
        model: Optional[str]: The audio model to use for transcription. Defaults to 'whisper-1'.
    """
    print(f'[Transcribing audio file from {file_path}]')
    try:
        # TODO: record the COST of the API call
        with open(file_path, 'rb') as audio_file:
            transcript = client.audio.translations.create(model=model, file=audio_file)
        print(transcript.text)

    except Exception as e:
        print(f'Error transcribing audio file: {e}')

@update_pwd_decorator
def parse_image(
    file_path: str, task: str = 'Describe this image as detail as possible.'
) -> None:
    """
    Parses the content of an image file and prints the description.

    Args:
        file_path: str: The path to the file to open.
        task: Optional[str]: The task description for the API call. Defaults to 'Describe this image as detail as possible.'.
    """
    print(f'[Reading image file from {file_path}]')
    # TODO: record the COST of the API call
    try:
        base64_image = _base64_img(file_path)
        response = client.chat.completions.create(
            model=OPENAI_MODEL,
            messages=_prepare_image_messages(task, base64_image),
            max_tokens=MAX_TOKEN,
        )
        content = response.choices[0].message.content
        print(content)

    except Exception as error:
        print(f'Error with the request: {error}')

@update_pwd_decorator
def parse_video(
    file_path: str,
    task: str = 'Describe this image as detail as possible.',
    frame_interval: int = 30,
) -> None:
    """
    Parses the content of an image file and prints the description.

    Args:
        file_path: str: The path to the video file to open.
        task: Optional[str]: The task description for the API call. Defaults to 'Describe this image as detail as possible.'.
        frame_interval: Optional[int]: The interval between frames to analyze. Defaults to 30.

    """
    print(
        f'[Processing video file from {file_path} with frame interval {frame_interval}]'
    )

    task = task or 'This is one frame from a video, please summarize this frame.'
    base64_frames = _base64_video(file_path)
    selected_frames = base64_frames[::frame_interval]

    if len(selected_frames) > 30:
        new_interval = len(base64_frames) // 30
        selected_frames = base64_frames[::new_interval]

    print(f'Totally {len(selected_frames)} would be analyze...\n')

    idx = 0
    for base64_frame in selected_frames:
        idx += 1
        print(f'Process the {file_path}, current No. {idx * frame_interval} frame...')
        # TODO: record the COST of the API call
        try:
            response = client.chat.completions.create(
                model=OPENAI_MODEL,
                messages=_prepare_image_messages(task, base64_frame),
                max_tokens=MAX_TOKEN,
            )

            content = response.choices[0].message.content
            current_frame_content = f"Frame {idx}'s content: {content}\n"
            print(current_frame_content)

        except Exception as error:
            print(f'Error with the request: {error}')

@update_pwd_decorator
def parse_pptx(file_path: str) -> None:
    """
    Parses the content of a pptx file and prints it.

    Args:
        file_path: str: The path to the file to open.
    """
    print(f'[Reading PowerPoint file from {file_path}]')
    try:
        pres = Presentation(str(file_path))
        text = []
        for slide_idx, slide in enumerate(pres.slides):
            text.append(f'@@ Slide {slide_idx + 1} @@')
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    text.append(shape.text)
        print('\n'.join(text))

    except Exception as e:
        print(f'Error reading PowerPoint file: {e}')

### MK Operations ###

@update_pwd_decorator
def get_mouse_position():
    """
    Get the current mouse position using xdotool.
    Returns:
        tuple: (x, y) coordinates of the mouse.
    """
    result = subprocess.run("xdotool getmouselocation --shell", shell=True, stdout=subprocess.PIPE)
    output = result.stdout.decode("utf-8")
    # Parse the output to extract x and y positions
    lines = output.strip().split("\n")
    x = int(lines[0].split("=")[1])  # Extract x-coordinate
    y = int(lines[1].split("=")[1])  # Extract y-coordinate
    return x, y

# @update_pwd_decorator
# def draw_dot(last_top_left, last_length, coordination):
#     """
#     Draws a red dot on the image at the specified coordination within a cropped screenshot.

#     Args:
#         last_top_left (tuple): Top-left corner of the region to crop from the screenshot (x, y).
#         last_length (int): Length of the region to crop (width is calculated as 3/4 of the length).
#         coordination (tuple): The (x, y) coordinates of the dot within the cropped region.

#     Returns:
#         Image: The modified image with the dot drawn.
#     """
#     # Capture a screenshot
#     screenshot = ImageGrab.grab()

#     # Convert dimensions to integers
#     last_length = int(last_length)

#     # Calculate dimensions for cropping
#     last_width = int(last_length * 5 // 8)  # Width of the cropped rectangle

#     # Validate if coordination is within the cropped region
#     x, y = last_top_left
#     dot_x, dot_y = coordination
#     if not (x <= dot_x <= x + last_length and y <= dot_y <= y + last_width):
#         raise ValueError(
#             f"Error: The dot at coordination {coordination} is outside the cropped region "
#             f"({last_length}x{last_width})."
#         )

#     # Draw a red dot on the original screenshot
#     draw = ImageDraw.Draw(screenshot)
#     dot_radius = max(3, int(min(last_width, last_length) * 0.02))  # Adjust dot radius based on cropped size
#     dot_color = (200, 0, 0)  # RGB color for the dot (red)
#     bbox = [
#         dot_x - dot_radius, dot_y - dot_radius,
#         dot_x + dot_radius, dot_y + dot_radius
#     ]

#     draw.ellipse(bbox, fill=dot_color, outline=dot_color)
    
#     # Crop the region defined by last_top_left and last_length
#     x, y = last_top_left
#     cropped_img = screenshot.crop((x, y, x + last_length, y + last_width))
    
#     # Save the screenshot with the drawn dot
#     screenshot_dir = "/workspace/screenshots"
#     os.makedirs(screenshot_dir, exist_ok=True)
#     timestamp = int(time.time())
#     screenshot_path = f"{screenshot_dir}/{timestamp}_draw_dot.png"
#     cropped_img.save(screenshot_path)
#     time.sleep(2)
#     print(f"<Screenshot saved at> {screenshot_path}")
     
# @update_pwd_decorator
# def draw_rectangle(last_top_left, last_length, new_top_left, new_length):
#     """
#     1. Captures a screenshot and crops a rectangular region defined by last_top_left and last_length.
#     2. Draws a rectangle inside the cropped region, specified by new_top_left and new_length.
#     3. Validates if new_top_left and new_length fall within the cropped region. If not, prints an error message and exits.
#     4. Saves the screenshot and modified image to a specified directory.

#     Args:
#         last_top_left (tuple): Top-left corner of the region to crop from the screenshot (x, y).
#         last_length (int): Length of the rectangle to crop (width is calculated as 5/8 of the length).
#         new_top_left (tuple): Top-left corner of the rectangle to draw within the cropped region (x, y).
#         new_length (int): Length of the rectangle to draw (width is calculated as 5/8 of the length).
#         command (str): Command or description for labeling the screenshot file.

#     Returns:
#         Image: The modified image with the drawn rectangle.
#     """
#     # Capture a screenshot
#     screenshot = ImageGrab.grab()

#     # Convert dimensions to integers
#     last_length = int(last_length)
#     new_length = int(new_length)

#     # Calculate dimensions for cropping and drawing
#     last_width = int(last_length * 5 // 8)  # Width of the cropped rectangle
#     new_width = int(new_length * 5 // 8)   # Width of the rectangle to draw

#     # Draw the rectangle on the screenshot
#     rect_color = (200, 0, 0)
#     if last_length >= 960:
#         rect_width = 10
#     elif 960 > last_length >= 240:
#         rect_width = 6
#     else:
#         rect_width = 3
#     draw = ImageDraw.Draw(screenshot)
#     rect_x, rect_y = new_top_left
#     rect_bbox = [
#         rect_x, rect_y,
#         rect_x + new_length, rect_y + new_width
#     ]
#     draw.rectangle(rect_bbox, outline=rect_color, width=rect_width)

#     # Crop the region defined by last_top_left and last_length
#     x, y = last_top_left
#     cropped_img = screenshot.crop((x, y, x + last_length, y + last_width))

#     # Validate if new_top_left and new_length are within the cropped region
#     nx, ny = new_top_left
#     if not (x <= nx <= x + last_length and y <= ny <= y + last_width and
#             x <= nx + new_length <= x + last_length and y <= ny + new_width <= y + last_width):
#         print("Error: The rectangle defined by top_left and length is outside the cropped region."
#               "Please reselect the top_left and length.")

#     # Save the screenshot and labeled image
#     screenshot_dir = "/workspace/screenshots"
#     os.makedirs(screenshot_dir, exist_ok=True)
#     timestamp = int(time.time())
    
#     screenshot_path = f"{screenshot_dir}/{timestamp}_draw_rectangle.png"

#     # Save original screenshot and labeled image
#     cropped_img.save(screenshot_path)
#     time.sleep(2)
#     print(f"<Screenshot saved at> {screenshot_path}")

# def draw_grid(img, length, offset=(0, 0)):
#     """
#     Draws a grid with fixed divisions (2 horizontal lines and 3 vertical lines),
#     and labels the intersection points based on global coordinates.
#     Now, adds a black outer frame region as a background, and places red coordinate labels
#     on the frame region instead of inside the grid.
#     """
#     from PIL import Image, ImageDraw, ImageFont

#     draw = ImageDraw.Draw(img)
#     width, height = img.size
#     if offset:
#         offset_x, offset_y = offset
#     else:
#         offset_x, offset_y = 0, 0

#     # Fixed divisions: 3 horizontal regions and 4 vertical regions
#     num_horizontal_lines = 5
#     num_vertical_lines = 7

#     # Calculate step sizes based on the screen size
#     step_x = width // (num_vertical_lines + 1)
#     step_y = height // (num_horizontal_lines + 1)

#     # Define line color, line style, and font size
#     grid_color = (200, 0, 0)  # Deep red (RGB)
#     label_color = "red"  # Red labels
#     dash_length = 10  # Length of each dash for dashed lines
#     font_size = max(int((length / 40)), 
#                     int(sqrt((width * height) / 2560)),
#                     8)  # Adjust font size based on screen area

#     # Load font
#     try:
#         font = ImageFont.truetype("DejaVuSans.ttf", font_size)
#     except:
#         font = ImageFont.load_default()

#     # Use textbbox to calculate maximum label dimensions
#     max_x_label_width = 0
#     max_y_label_width = 0

#     for i in range(num_vertical_lines + 2):
#         label = f"{i * step_x + offset_x}"
#         bbox = draw.textbbox((0, 0), label, font=font)
#         max_x_label_width = max(max_x_label_width, bbox[2] - bbox[0])

#     for j in range(num_horizontal_lines + 2):
#         label = f"{j * step_y + offset_y}"
#         bbox = draw.textbbox((0, 0), label, font=font)
#         max_y_label_width = max(max_y_label_width, bbox[2] - bbox[0])

#     frame_thickness_x = max_x_label_width  # Match longest X-coordinate label
#     frame_thickness_y = font_size  # Match font size for Y-coordinate label height

#     # print(f"Frame thickness: {frame_thickness_x} x {frame_thickness_y}")

#     # Create a new image with the black frame
#     new_width = width + 2 * frame_thickness_x
#     new_height = height + 2 * frame_thickness_y
#     new_img = Image.new("RGB", (new_width, new_height), "black")

#     # Paste the original image onto the new background
#     new_img.paste(img, (frame_thickness_x, frame_thickness_y))
#     draw = ImageDraw.Draw(new_img)

#     # Draw vertical dashed lines
#     for i in range(1, num_vertical_lines + 1):
#         x = i * step_x + frame_thickness_x
#         for y in range(frame_thickness_y, height + frame_thickness_y, dash_length * 2):
#             draw.line([(x, y), (x, y + dash_length)], fill=grid_color, width=1)

#     # Draw horizontal dashed lines
#     for i in range(1, num_horizontal_lines + 1):
#         y = i * step_y + frame_thickness_y
#         for x in range(frame_thickness_x, width + frame_thickness_x, dash_length * 2):
#             draw.line([(x, y), (x + dash_length, y)], fill=grid_color, width=1)

#     # Label X-coordinates on the top and bottom frame
#     for i in range(num_vertical_lines + 2):
#         x = i * step_x + frame_thickness_x
#         label = f"{i * step_x + offset_x}"
#         # Top frame
#         draw.text(
#             (x - max_x_label_width // 2, frame_thickness_y // 2 - font_size // 2), label, fill=label_color, font=font
#         )
#         # Bottom frame
#         draw.text(
#             (x - max_x_label_width // 2, height + frame_thickness_y), label, fill=label_color, font=font
#         )

#     # Label Y-coordinates on the left and right frame
#     for j in range(num_horizontal_lines + 2):
#         y = j * step_y + frame_thickness_y
#         label = f"{j * step_y + offset_y}"
#         # Left frame
#         draw.text(
#             (frame_thickness_x // 2 - max_y_label_width // 2, y - font_size // 2), label, fill=label_color, font=font
#         )
#         # Right frame
#         draw.text(
#             (width + frame_thickness_x, y - font_size // 2), label, fill=label_color, font=font
#         )

#     return new_img

@update_pwd_decorator
def take_screenshot(command: str | None = None, top_left: tuple | None = None, length: int | None = None) -> str:
    """
    Captures a screenshot, adds a coordinate system with dashed deep red lines, larger labels,
    and marks the mouse position with an enlarged red dot. Each intersection of the coordinate
    grid is labeled with its (x, y) position.
    top_left & length: tuple: (x, y) coordinates of the region to capture and length of the side.
    """
    # Constants
    screenshot_dir = "/workspace/screenshots"
    os.makedirs(screenshot_dir, exist_ok=True)
    timestamp = int(time.time())
    if command:
        # safe_command = re.sub(r'[^\w\s-]', '_', command)  # FIXME: iGnore this for now
        # safe_command = re.sub(r'\s+', '_', safe_command)  
        screenshot_path = f"{screenshot_dir}/{timestamp}.png"
        # labeled_screenshot_path = f"{screenshot_dir}/{timestamp}_label.png"
    else:
        screenshot_path = f"{screenshot_dir}/{timestamp}_screenshot.png"
        # labeled_screenshot_path = f"{screenshot_dir}/{timestamp}_screenshot_label.png"

    # Get screen resolution
    screen_width, screen_height = ImageGrab.grab().size

    # Define the capture region
    if top_left and length:
        # print(f"top_left: {top_left}, length: {length}")
        region = (
            int(top_left[0]),
            int(top_left[1]),
            min(int(top_left[0]) + int(length), screen_width),
            min(int(top_left[1]) + int(length * 0.75), screen_height),
        )
        img = ImageGrab.grab(bbox=region)
        # print(f"Capturing region: {region}")
        img.save(screenshot_path)
        time.sleep(2)
        print(f"<Screenshot saved at> {screenshot_path}")
        length = region[2] - region[0]
    else:
        img = ImageGrab.grab()
        img.save(screenshot_path)
        time.sleep(2)
        print(f"<Screenshot saved at> {screenshot_path}")
        length = screen_width
        
    # new_img = draw_grid(img, length=length, offset=top_left)

    # # Save the labeled image
    # new_img.save(labeled_screenshot_path)
    # time.sleep(2)
    # print(f"<Screenshot saved at> {labeled_screenshot_path}")

@update_pwd_decorator
def mouse_left_click(x, y, button="left"):
    """
    Simulates a mouse click at the specified coordinates.
    """
    if x != -1 and y != -1:
        button_code = 1 if button == "left" else 3
        subprocess.run(f"xdotool mousemove {x} {y} click {button_code}", shell=True)
        time.sleep(1)
        take_screenshot(f'mouse_left_click({x}, {y})')
    else:
        print('Please provide a more detailed description of where you want to click.')

@update_pwd_decorator
def mouse_double_click(x, y, button="left"):
    """
    Simulates a double-click at the specified coordinates.
    """
    if x != -1 and y != -1:
        button_code = 1 if button == "left" else 3
        subprocess.run(f"xdotool mousemove {x} {y} click {button_code} click {button_code}", shell=True)
        time.sleep(1)
        take_screenshot(f'mouse_double_click({x}, {y})')
    else:
        print('Please provide a more detailed description of where you want to double-click.')

@update_pwd_decorator
def mouse_right_click(x, y):
    """
    Simulates a right mouse click at the specified coordinates.
    :param x: x-coordinate.
    :param y: y-coordinate.
    """
    if x != -1 and y != -1:
        subprocess.run(f"xdotool mousemove {x} {y} click 3", shell=True)
        time.sleep(1)
        take_screenshot(f'mouse_right_click({x}, {y})')
    else:
        print('Please provide a more detailed description of where you want to right-click.')

@update_pwd_decorator
def mouse_move(x, y):
    """
    Moves the mouse to the specified coordinates.
    """
    if x != -1 and y != -1:
        subprocess.run(f"xdotool mousemove {x} {y}", shell=True)
        time.sleep(1)
        take_screenshot(f'mouse_move({x}, {y})')
    else:
        print('Please provide a more detailed description of where you want to move the mouse.')

@update_pwd_decorator
def mouse_scroll(direction="up", amount=1):
    """
    Simulates a mouse scroll up or down.
    """
    button_code = 4 if direction == "up" else 5
    for _ in range(amount):
        subprocess.run(f"xdotool click {button_code}", shell=True)
        time.sleep(1)
    take_screenshot(f'mouse_scroll({direction}, {amount})')

@update_pwd_decorator
def type_text(text: str):
    """
    Simulates typing text.
    """
    subprocess.run(f'xdotool type "{text}"', shell=True)
    sanitized_text = text.replace('/', '_')
    time.sleep(1)
    take_screenshot(f'type_text({sanitized_text})')

@update_pwd_decorator
def press_key(key: str):
    """
    Simulates pressing a specific key.
    """
    subprocess.run(f"xdotool key {key}", shell=True)
    sanitized_key = key.replace('/', '_')
    time.sleep(1)
    take_screenshot(f'press_key({sanitized_key})')

@update_pwd_decorator
def open_application(app_name):
    """
    Opens a specific application using xdotool.
    :param app_name: The name of the application to open (e.g., "chrome", "gedit").
    """
    # Simulate pressing the "Super" key to open the application launcher
    subprocess.run("xdotool key super", shell=True)
    time.sleep(1)
    # Type the application name in the search bar
    subprocess.run(f"xdotool type '{app_name}'", shell=True)
    time.sleep(1)
    # Press "Enter" to open the application
    subprocess.run("xdotool key Return", shell=True)
    print(f"Opening {app_name}...")
    time.sleep(2)
    take_screenshot(f'open_application({app_name})')

@update_pwd_decorator
def mouse_drag(x_start, y_start, x_end, y_end, button="left"):
    """
    Simulates dragging the mouse from one position to another with smooth movements.
    :param x_start: Starting x-coordinate.
    :param y_start: Starting y-coordinate.
    :param x_end: Ending x-coordinate.
    :param y_end: Ending y-coordinate.
    :param button: Mouse button to hold during the drag ("left" or "right").
    :param delay: Delay (in seconds) between each step.
    
    Mouse drag may not be very accurate in Linux as the screen is sperated into multiple cells.
    """
    
    try:
        button_code = 1 if button == "left" else 3
        
        # Move to start position
        subprocess.run(f"xdotool mousemove {x_start} {y_start}", shell=True, check=True)
        time.sleep(0.2)
        
        # Press the mouse button
        subprocess.run(f"xdotool mousedown {button_code}", shell=True, check=True)
        time.sleep(0.2)
        
        mouse_pos_x, mouse_pos_y = get_mouse_position()
        print(f"Before drag: Mouse position: {mouse_pos_x}, {mouse_pos_y}")
        
        # Compute incremental steps for smooth movement
        # Move horizontally first
        x_step = x_end - x_start
        for i in range(1, abs(x_step) + 1):
            if x_step > 0:
                intermediate_x = int(x_start + i)
            else:
                intermediate_x = int(x_start - i)
            subprocess.run(f"xdotool mousemove {intermediate_x} {y_start}", shell=True, check=True)
        
        # Move vertically next
        y_step = y_end - y_start
        for i in range(1, abs(y_step) + 1):
            if y_step > 0:
                intermediate_y = int(y_start + i)
            else:
                intermediate_y = int(y_start - i)
            subprocess.run(f"xdotool mousemove {x_end} {intermediate_y}", shell=True, check=True)
        
        # Release the mouse button
        subprocess.run(f"xdotool mouseup {button_code}", shell=True, check=True)
        time.sleep(0.2)
        
        mouse_pos_x, mouse_pos_y = get_mouse_position()
        print(f"After drag: Mouse position: {mouse_pos_x}, {mouse_pos_y}")
        
        # Optionally capture a screenshot
        take_screenshot(f'mouse_drag({x_start}, {y_start}, {x_end}, {y_end})')
    except subprocess.CalledProcessError as e:
        print(f"Error executing xdotool command: {e}")
    except Exception as e:
        print(f"Unexpected error: {e}")

@update_pwd_decorator
def mouse_right_click(x, y):
    """
    Simulates a right mouse click at the specified coordinates.
    :param x: x-coordinate.
    :param y: y-coordinate.
    """
    subprocess.run(f"xdotool mousemove {x} {y} click 3", shell=True)
    time.sleep(1)
    take_screenshot(f'mouse_right_click({x}, {y})')

@update_pwd_decorator
def localization(top_left: tuple | None = None, length: int | None = None):
    """
    Localizes the mouse position by capturing a screenshot and adding a coordinate system.

    Args:
        top_left (tuple | None): 
            The top-left corner of the screenshot region as a tuple of (x, y) coordinates. 
            If None, the screenshot will cover the entire screen. Defaults to None.
        
        length (int | None): 
            The side length of the screenshot region, forming a square. 
            If None, the screenshot region will cover the entire screen or be determined dynamically. Defaults to None.

    Returns:
        None
    """
    print("Localizing mouse position...")
    take_screenshot("localization", top_left, length)

__all__ = [
    # file operation
    'open_file',
    'goto_line',
    'scroll_down',
    'scroll_up',
    'create_file',
    'append_file',
    'edit_file',
    'search_dir',
    'search_file',
    'find_file',
    'replace_function',
    'search_function',
    # readers
    'parse_pdf',
    'parse_docx',
    'parse_latex',
    'parse_pptx',
    # tools
    'parse_audio',
    'parse_video',
    'parse_image',
    # MK operations
    'take_screenshot',
    'mouse_left_click',
    'mouse_right_click',
    'mouse_double_click',
    'mouse_move',
    'mouse_scroll',
    'type_text',
    'press_key',
    'open_application',
    'mouse_drag',
    'mouse_right_click',
    'localization',
]

if OPENAI_API_KEY and OPENAI_BASE_URL:
    __all__ += ['parse_audio', 'parse_video', 'parse_image']

DOCUMENTATION = ''
for func_name in __all__:
    func = globals()[func_name]

    cur_doc = func.__doc__
    # remove indentation from docstring and extra empty lines
    cur_doc = '\n'.join(filter(None, map(lambda x: x.strip(), cur_doc.split('\n'))))
    # now add a consistent 4 indentation
    cur_doc = '\n'.join(map(lambda x: ' ' * 4 + x, cur_doc.split('\n')))

    fn_signature = f'{func.__name__}' + str(signature(func))
    DOCUMENTATION += f'{fn_signature}:\n{cur_doc}\n\n'
